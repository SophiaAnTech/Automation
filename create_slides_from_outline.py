from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TXT_FILE    = "outline.txt"
OUTPUT_FILE = "sales_overview.pptx"


def parse_txt(filepath):
    slides = []
    current_title = None
    current_bullets = []

    with open(filepath, "r", encoding="utf-8") as f:
        lines = f.readlines()

    for line in lines:
        stripped = line.strip()

        if stripped == "":
            continue

        # original line doesn't start with 3 spaces, it is at the top indentation level and 
        if stripped.startswith("* ") and not line.startswith("   "):
            # the current line is for the next slide, save the previous slide first
            if current_title is not None:
                slides.append({"title": current_title, "bullets": current_bullets})
            current_title = stripped[2:]
            current_bullets = []

        elif stripped.startswith("* "):
            current_bullets.append(stripped[2:])

    # save the last slide
    if current_title is not None:
        slides.append({"title": current_title, "bullets": current_bullets})

    return slides


def create_presentation(slides):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    bullet_layout = prs.slide_layouts[1]  # Title and Content layout

    for slide_data in slides:
        slide = prs.slides.add_slide(bullet_layout)

        title_box   = slide.shapes.title
        content_box = slide.placeholders[1]

        title_box.text = slide_data["title"]
        title_box.text_frame.paragraphs[0].runs[0].font.size  = Pt(36)
        title_box.text_frame.paragraphs[0].runs[0].font.bold  = True
        title_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x27, 0x61)

        tf = content_box.text_frame
        tf.clear()

        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            para.text = bullet
            para.level = 0
            run = para.runs[0]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT_FILE)
    print("Saved to " + OUTPUT_FILE)


slides = parse_txt(TXT_FILE)
create_presentation(slides)
